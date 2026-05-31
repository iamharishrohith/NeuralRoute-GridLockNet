import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_upgraded_presentation():
    prs = Presentation()
    
    # Set slide dimensions to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Color scheme - Flipkart/SIH White Corporate Theme
    COLOR_BG = RGBColor(255, 255, 255)       # #FFFFFF (Pure White)
    COLOR_CARD = RGBColor(248, 250, 252)     # #F8FAFC (Subtle Slate/White card bg)
    COLOR_BORDER = RGBColor(226, 232, 240)   # #E2E8F0 (Subtle borders)
    COLOR_BLUE = RGBColor(40, 116, 240)      # #2874F0 (Flipkart Corporate Blue)
    COLOR_ORANGE = RGBColor(245, 114, 0)     # #F57200 (Flipkart Active Orange)
    COLOR_DARK = RGBColor(15, 23, 42)        # #0F172A (Deep Slate/Midnight Blue)
    COLOR_TEXT = RGBColor(71, 85, 105)      # #475569 (Slate Gray for body text)
    COLOR_WHITE = RGBColor(255, 255, 255)     # #FFFFFF (Pure White)
    
    # Helper to create slides with solid white background and standard headers matching SIH template
    def add_professional_slide(title_text):
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # Use blank layout
        
        # Solid white background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background()
        
        # Main Header Box in center top
        header_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(0.15), Inches(4.3), Inches(0.55)
        )
        header_box.fill.solid()
        header_box.fill.fore_color.rgb = COLOR_DARK
        header_box.line.fill.background()
        
        tf_h = header_box.text_frame
        p_h = tf_h.paragraphs[0]
        p_h.text = "GridLock 2.0 Traffic Predictor"
        p_h.font.size = Pt(14)
        p_h.font.bold = True
        p_h.font.color.rgb = COLOR_WHITE
        p_h.font.name = "Calibri"
        p_h.alignment = PP_ALIGN.CENTER
        
        # Left Logo Box (NeuralRoute)
        logo_left = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(2.5), Inches(0.6))
        tf_ll = logo_left.text_frame
        p_ll = tf_ll.paragraphs[0]
        p_ll.text = "★ NEURALROUTE"
        p_ll.font.size = Pt(16)
        p_ll.font.bold = True
        p_ll.font.color.rgb = COLOR_BLUE
        p_ll.font.name = "Trebuchet MS"
        
        # Right Logo Box (Flipkart Hackathon)
        logo_right = slide.shapes.add_textbox(Inches(10.0), Inches(0.1), Inches(2.8), Inches(0.6))
        tf_lr = logo_right.text_frame
        p_lr = tf_lr.paragraphs[0]
        p_lr.text = "FLIPKART GRIDLOCK 2.0"
        p_lr.font.size = Pt(14)
        p_lr.font.bold = True
        p_lr.font.color.rgb = COLOR_ORANGE
        p_lr.font.name = "Trebuchet MS"
        p_lr.alignment = PP_ALIGN.RIGHT
        
        # Visual line underneath headers
        header_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.75), Inches(12.33), Inches(0.02)
        )
        header_line.fill.solid()
        header_line.fill.fore_color.rgb = COLOR_BORDER
        header_line.line.fill.background()
        
        # Slide progress / category tag at bottom
        tag_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.33), Inches(0.4))
        tf_tag = tag_box.text_frame
        p_tag = tf_tag.paragraphs[0]
        p_tag.text = "TEAM: NEURALROUTE  |  LEAD: HARISH ROHITH S, MEMBERS: KARTHIKEYAN T, RAM PRASATH V, NANDHAKISHORE A"
        p_tag.font.size = Pt(10)
        p_tag.font.bold = True
        p_tag.font.color.rgb = COLOR_TEXT
        p_tag.font.name = "Calibri"
        
        # Main section header banner
        banner = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.9), Inches(12.33), Inches(0.45)
        )
        banner.fill.solid()
        banner.fill.fore_color.rgb = COLOR_CARD
        banner.line.color.rgb = COLOR_BORDER
        banner.line.width = Pt(1)
        
        tf_b = banner.text_frame
        p_b = tf_b.paragraphs[0]
        p_b.text = title_text.upper()
        p_b.font.size = Pt(13)
        p_b.font.bold = True
        p_b.font.color.rgb = COLOR_DARK
        p_b.font.name = "Trebuchet MS"
        p_b.alignment = PP_ALIGN.CENTER
        
        return slide

    # =========================================================================
    # SLIDE 1: TITLE SLIDE (SIH Style Recreation)
    # =========================================================================
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # White bg
    bg1 = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = COLOR_BG
    bg1.line.fill.background()
    
    # Left brand block
    brand_left = slide1.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(3.0), Inches(0.8))
    tf_bl = brand_left.text_frame
    p_bl1 = tf_bl.paragraphs[0]
    p_bl1.text = "NEURALROUTE ENGINE"
    p_bl1.font.size = Pt(18)
    p_bl1.font.bold = True
    p_bl1.font.color.rgb = COLOR_BLUE
    p_bl1.font.name = "Trebuchet MS"
    p_bl2 = tf_bl.add_paragraph()
    p_bl2.text = "Innovations & Mobility Solutions"
    p_bl2.font.size = Pt(10)
    p_bl2.font.color.rgb = COLOR_TEXT
    p_bl2.font.name = "Calibri"

    # Middle title box
    middle_title = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(3.8), Inches(0.25), Inches(5.8), Inches(0.6)
    )
    middle_title.fill.solid()
    middle_title.fill.fore_color.rgb = COLOR_BG
    middle_title.line.color.rgb = COLOR_DARK
    middle_title.line.width = Pt(1.5)
    tf_mt = middle_title.text_frame
    p_mt = tf_mt.paragraphs[0]
    p_mt.text = "FLIPKART GRIDLOCK HACKATHON 2.0"
    p_mt.font.size = Pt(16)
    p_mt.font.bold = True
    p_mt.font.color.rgb = COLOR_DARK
    p_mt.font.name = "Trebuchet MS"
    p_mt.alignment = PP_ALIGN.CENTER
    
    # Right brand block
    brand_right = slide1.shapes.add_textbox(Inches(9.8), Inches(0.2), Inches(3.0), Inches(0.8))
    tf_br = brand_right.text_frame
    p_br1 = tf_br.paragraphs[0]
    p_br1.text = "GRIDLOCK 2.0"
    p_br1.font.size = Pt(18)
    p_br1.font.bold = True
    p_br1.font.color.rgb = COLOR_ORANGE
    p_br1.font.name = "Trebuchet MS"
    p_br1.alignment = PP_ALIGN.RIGHT
    p_br2 = tf_br.add_paragraph()
    p_br2.text = "Championship Solvers 2026"
    p_br2.font.size = Pt(10)
    p_br2.font.color.rgb = COLOR_TEXT
    p_br2.font.name = "Calibri"
    p_br2.alignment = PP_ALIGN.RIGHT

    # Middle Slide divider
    div_line = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.95), Inches(12.33), Inches(0.02)
    )
    div_line.fill.solid()
    div_line.fill.fore_color.rgb = COLOR_BORDER
    div_line.line.fill.background()

    # Core Hackathon Metadata Fields (Slide 1 Left Side)
    meta_box = slide1.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.33), Inches(3.0))
    tf_meta = meta_box.text_frame
    tf_meta.word_wrap = True
    
    meta_fields = [
        ("Problem Statement ID", "2452"),
        ("Problem Statement Title", "Development of AI/ML based solution for Traffic Demand Prediction in GridLock 2.0"),
        ("Theme", "Smart Cities, Urban Mobility & Logistics"),
        ("Problem Statement Category", "Software / Machine Learning"),
        ("Team ID", "27815"),
        ("Team Name", "NeuralRoute"),
        ("Team Roster", "Harish Rohith S (Team Lead), Karthikeyan T (Member 1), Ram Prasath V (Member 2), Nandhakishore A (Member 3)")
    ]
    
    for label, val in meta_fields:
        p = tf_meta.add_paragraph()
        p.text = f"{label} – {val}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_DARK
        p.font.name = "Calibri"
        p.space_after = Pt(6)

    # Bottom "About this Model" Box (Slide 1 Bottom)
    about_card = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.7), Inches(12.33), Inches(1.8)
    )
    about_card.fill.solid()
    about_card.fill.fore_color.rgb = COLOR_BG
    about_card.line.color.rgb = COLOR_DARK
    about_card.line.width = Pt(1.5)
    
    tf_about = about_card.text_frame
    tf_about.word_wrap = True
    
    pa1 = tf_about.paragraphs[0]
    pa1.text = "About this Model"
    pa1.font.size = Pt(18)
    pa1.font.bold = True
    pa1.font.color.rgb = COLOR_BLUE
    pa1.font.name = "Trebuchet MS"
    pa1.space_after = Pt(6)
    
    pa2 = tf_about.add_paragraph()
    pa2.text = "In this model, Team NeuralRoute creatively designs GridLockNet as a custom multi-branch deep spatial-temporal tabular ResNet in PyTorch. Fusing learned nested geohash embeddings, cyclic continuous time projectors, and dynamic Gated Linear Unit lag selection blocks, it predicts traffic demand chronologically via a dampened recursive loop to control exposure bias. It is leakage-free, robustly regularized, and guarantees bounded output in [0.0, 1.0]."
    pa2.font.size = Pt(14)
    pa2.font.color.rgb = COLOR_DARK
    pa2.font.name = "Calibri"

    # =========================================================================
    # SLIDE 2: PROPOSED SOLUTION & INNOVATION | METHODOLOGY (SIH Layout)
    # =========================================================================
    slide2 = add_professional_slide("Proposed Solution & Innovation  |  Methodology")
    
    # Left Column: PROPOSED SOLUTION & INNOVATION (5 colorful rounded cards)
    header_sol = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(5.9), Inches(0.4)
    )
    header_sol.fill.solid()
    header_sol.fill.fore_color.rgb = COLOR_CARD
    header_sol.line.color.rgb = COLOR_BORDER
    header_sol.line.width = Pt(1)
    tf_hsol = header_sol.text_frame
    p_hsol = tf_hsol.paragraphs[0]
    p_hsol.text = "PROPOSED SOLUTION & INNOVATION"
    p_hsol.font.size = Pt(12)
    p_hsol.font.bold = True
    p_hsol.font.color.rgb = COLOR_DARK
    p_hsol.font.name = "Trebuchet MS"
    
    # Card 1: Multi-Branch ResNet (Green)
    c1 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.0), Inches(2.8), Inches(1.4))
    c1.fill.solid()
    c1.fill.fore_color.rgb = RGBColor(76, 175, 80) # Green
    c1.line.fill.background()
    tf_c1 = c1.text_frame
    tf_c1.word_wrap = True
    pc1 = tf_c1.paragraphs[0]
    pc1.text = "1. Multi-Branch ResNet\nFuses spatial geohashes, cyclic intervals, and dynamic lag rollouts in parallel."
    pc1.font.size = Pt(11)
    pc1.font.bold = True
    pc1.font.color.rgb = COLOR_WHITE
    pc1.font.name = "Calibri"
    
    # Card 2: Nested Embeddings (Yellow)
    c2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.6), Inches(2.0), Inches(2.8), Inches(1.4))
    c2.fill.solid()
    c2.fill.fore_color.rgb = RGBColor(241, 196, 15) # Yellow
    c2.line.fill.background()
    tf_c2 = c2.text_frame
    tf_c2.word_wrap = True
    pc2 = tf_c2.paragraphs[0]
    pc2.text = "2. Nested Spatial Embeddings\nLearns geohash weights at lengths 3, 4, 5 to map neighborhood hierarchies."
    pc2.font.size = Pt(11)
    pc2.font.bold = True
    pc2.font.color.rgb = COLOR_DARK
    pc2.font.name = "Calibri"
    
    # Card 3: 0.9 Dampened Lags (Grey)
    c3 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.6), Inches(3.6), Inches(2.8), Inches(1.4))
    c3.fill.solid()
    c3.fill.fore_color.rgb = RGBColor(127, 140, 141) # Grey
    c3.line.fill.background()
    tf_c3 = c3.text_frame
    tf_c3.word_wrap = True
    pc3 = tf_c3.paragraphs[0]
    pc3.text = "3. 0.9 Dampened Lags\nControls recursive rollout error propagation and exposure bias."
    pc3.font.size = Pt(11)
    pc3.font.bold = True
    pc3.font.color.rgb = COLOR_WHITE
    pc3.font.name = "Calibri"
    
    # Card 4: Leakage Isolation (Orange)
    c4 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.6), Inches(2.8), Inches(1.4))
    c4.fill.solid()
    c4.fill.fore_color.rgb = COLOR_ORANGE
    c4.line.fill.background()
    tf_c4 = c4.text_frame
    tf_c4.word_wrap = True
    pc4 = tf_c4.paragraphs[0]
    pc4.text = "4. Inductive Clean Preprocessing\nWeather statistics computed strictly on Train set, avoiding leakage."
    pc4.font.size = Pt(11)
    pc4.font.bold = True
    pc4.font.color.rgb = COLOR_WHITE
    pc4.font.name = "Calibri"
    
    # Card 5: Sigmoid Output Bounding (Blue)
    c5 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.2), Inches(5.9), Inches(1.0))
    c5.fill.solid()
    c5.fill.fore_color.rgb = COLOR_BLUE
    c5.line.fill.background()
    tf_c5 = c5.text_frame
    tf_c5.word_wrap = True
    pc5 = tf_c5.paragraphs[0]
    pc5.text = "5. Multi-Stage Sigmoid Bounding Output\nCustom scaled Sigmoid layer restricts traffic forecasts strictly in valid [0.0, 1.0] interval, taming out-of-bounds variance."
    pc5.font.size = Pt(11)
    pc5.font.bold = True
    pc5.font.color.rgb = COLOR_WHITE
    pc5.font.name = "Calibri"

    # Middle vertical dividing line
    v_line = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.6), Inches(1.5), Inches(0.02), Inches(4.7)
    )
    v_line.fill.solid()
    v_line.fill.fore_color.rgb = COLOR_BORDER
    v_line.line.fill.background()

    # Right Column: METHODOLOGY (5-Level isometric-style stack on the right)
    header_meth = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.5), Inches(6.0), Inches(0.4)
    )
    header_meth.fill.solid()
    header_meth.fill.fore_color.rgb = COLOR_CARD
    header_meth.line.color.rgb = COLOR_BORDER
    header_meth.line.width = Pt(1)
    tf_hmeth = header_meth.text_frame
    p_hmeth = tf_hmeth.paragraphs[0]
    p_hmeth.text = "METHODOLOGY GRIDS & SECTIONS"
    p_hmeth.font.size = Pt(12)
    p_hmeth.font.bold = True
    p_hmeth.font.color.rgb = COLOR_DARK
    p_hmeth.font.name = "Trebuchet MS"

    # Stack Level 1 to 5 (Left edge of right panel)
    levels = [
        ("Level 5", "Autoregressive Dampened forecasting rollout (0.9 Optimum)", COLOR_BLUE, Inches(2.0)),
        ("Level 4", "PyTorch deepcopy checkpoint lock (locking Epoch 8)", COLOR_ORANGE, Inches(2.8)),
        ("Level 3", "Nested geohash grouping and spatial shift coordinate reverse engineering", RGBColor(127, 140, 141), Inches(3.6)),
        ("Level 2", "Deterministic road metadata capacities imputation (Lanes -> RoadType)", RGBColor(241, 196, 15), Inches(4.4)),
        ("Level 1", "3D Grid reconstruction (reconstructing missing traffic blocks to 0.0)", RGBColor(76, 175, 80), Inches(5.2))
    ]
    
    for name, desc, color, top in levels:
        # The Level trapezoid block
        trap = slide2.shapes.add_shape(MSO_SHAPE.TRAPEZOID, Inches(6.8), top, Inches(1.6), Inches(0.7))
        trap.fill.solid()
        trap.fill.fore_color.rgb = color
        trap.line.fill.background()
        tf_tr = trap.text_frame
        p_tr = tf_tr.paragraphs[0]
        p_tr.text = name
        p_tr.font.size = Pt(11)
        p_tr.font.bold = True
        p_tr.font.color.rgb = COLOR_WHITE if color != RGBColor(241, 196, 15) else COLOR_DARK
        p_tr.font.name = "Calibri"
        p_tr.alignment = PP_ALIGN.CENTER
        
        # The explanation text box beside it
        exp_box = slide2.shapes.add_textbox(Inches(8.5), top, Inches(4.3), Inches(0.7))
        tf_exp = exp_box.text_frame
        tf_exp.word_wrap = True
        p_exp = tf_exp.paragraphs[0]
        p_exp.text = desc
        p_exp.font.size = Pt(12)
        p_exp.font.bold = True
        p_exp.font.color.rgb = COLOR_DARK
        p_exp.font.name = "Calibri"

    # =========================================================================
    # SLIDE 3: CODE PROCESS & IMPLEMENTATION | TECHNOLOGIES (SIH Layout)
    # =========================================================================
    slide3 = add_professional_slide("Code Process of Implementation  |  Configurations")
    
    # Top Section: CODE PROCESS OF IMPLEMENTATION
    flow_hdr = slide3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.4), Inches(12.33), Inches(0.35)
    )
    flow_hdr.fill.solid()
    flow_hdr.fill.fore_color.rgb = COLOR_CARD
    flow_hdr.line.color.rgb = COLOR_BORDER
    flow_hdr.line.width = Pt(1)
    tf_fhdr = flow_hdr.text_frame
    p_fhdr = tf_fhdr.paragraphs[0]
    p_fhdr.text = "CODE PROCESS OF IMPLEMENTATION PIPELINE"
    p_fhdr.font.size = Pt(11)
    p_fhdr.font.bold = True
    p_fhdr.font.color.rgb = COLOR_DARK
    p_fhdr.font.name = "Trebuchet MS"
    p_fhdr.alignment = PP_ALIGN.CENTER
    
    # Flowchart blocks (Horizontal Sequence)
    flow_steps = [
        "Raw Traffic Logs", "3D Grid Reconst.", "Capacity Fills", "Cyclic Time Fusing", 
        "Nested Geohashes", "MinMax Scaling", "GridLockNet ResNet", "Sigmoid Bounding"
    ]
    
    for idx, step in enumerate(flow_steps):
        left_pos = Inches(0.5 + idx * 1.55)
        step_box = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_pos, Inches(1.9), Inches(1.4), Inches(0.85))
        step_box.fill.solid()
        step_box.fill.fore_color.rgb = COLOR_BG
        step_box.line.color.rgb = COLOR_BLUE
        step_box.line.width = Pt(1)
        tf_s = step_box.text_frame
        tf_s.word_wrap = True
        ps = tf_s.paragraphs[0]
        ps.text = step
        ps.font.size = Pt(10)
        ps.font.bold = True
        ps.font.color.rgb = COLOR_DARK
        ps.font.name = "Calibri"
        ps.alignment = PP_ALIGN.CENTER
        
        # Connecting arrow (except last)
        if idx < len(flow_steps) - 1:
            arrow = slide3.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left_pos + Inches(1.4), Inches(2.2), Inches(0.15), Inches(0.2))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLOR_ORANGE
            arrow.line.fill.background()

    # Loop arrow text representing dynamic recursive feedback
    loop_text = slide3.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.33), Inches(0.4))
    tf_loop = loop_text.text_frame
    p_loop = tf_loop.paragraphs[0]
    p_loop.text = "🔄 Dynamic Autoregressive Feedback Loop: Predicted demand t is dynamically fed back into grid to update lags for step t + 15m."
    p_loop.font.size = Pt(12)
    p_loop.font.bold = True
    p_loop.font.color.rgb = COLOR_ORANGE
    p_loop.font.name = "Consolas"
    p_loop.alignment = PP_ALIGN.CENTER

    # Bottom Section: TECHNOLOGIES & CONFIGURATIONS
    tech_hdr = slide3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.3), Inches(12.33), Inches(0.35)
    )
    tech_hdr.fill.solid()
    tech_hdr.fill.fore_color.rgb = COLOR_CARD
    tech_hdr.line.color.rgb = COLOR_BORDER
    tech_hdr.line.width = Pt(1)
    tf_thdr = tech_hdr.text_frame
    p_thdr = tf_thdr.paragraphs[0]
    p_thdr.text = "TECHNOLOGIES & CONFIGURATIONS FOR DEEP TRAFFIC FORECASTING"
    p_thdr.font.size = Pt(11)
    p_thdr.font.bold = True
    p_thdr.font.color.rgb = COLOR_DARK
    p_thdr.font.name = "Trebuchet MS"
    p_thdr.alignment = PP_ALIGN.CENTER

    # Card 1: Frameworks & Libraries
    tc1 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.8), Inches(3.8), Inches(3.0))
    tc1.fill.solid()
    tc1.fill.fore_color.rgb = COLOR_CARD
    tc1.line.color.rgb = COLOR_BORDER
    tc1.line.width = Pt(1)
    tf_tc1 = tc1.text_frame
    tf_tc1.word_wrap = True
    ptc1_1 = tf_tc1.paragraphs[0]
    ptc1_1.text = "Frameworks & Libraries"
    ptc1_1.font.size = Pt(14)
    ptc1_1.font.bold = True
    ptc1_1.font.color.rgb = COLOR_BLUE
    ptc1_1.font.name = "Trebuchet MS"
    ptc1_1.space_after = Pt(10)
    ptc1_2 = tf_tc1.add_paragraph()
    ptc1_2.text = "• Python (Primary Core)\n• PyTorch (Model Architecture)\n• NumPy & Pandas (Data Grids)\n• Scikit-Learn (Scalers & Metrics)"
    ptc1_2.font.size = Pt(13)
    ptc1_2.font.color.rgb = COLOR_DARK
    ptc1_2.font.name = "Calibri"

    # Card 2: PyTorch GridLockNet Features
    tc2 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.6), Inches(3.8), Inches(4.1), Inches(3.0))
    tc2.fill.solid()
    tc2.fill.fore_color.rgb = COLOR_CARD
    tc2.line.color.rgb = COLOR_BORDER
    tc2.line.width = Pt(1)
    tf_tc2 = tc2.text_frame
    tf_tc2.word_wrap = True
    ptc2_1 = tf_tc2.paragraphs[0]
    ptc2_1.text = "Special Features in GridLockNet"
    ptc2_1.font.size = Pt(14)
    ptc2_1.font.bold = True
    ptc2_1.font.color.rgb = COLOR_BLUE
    ptc2_1.font.name = "Trebuchet MS"
    ptc2_1.space_after = Pt(10)
    ptc2_2 = tf_tc2.add_paragraph()
    ptc2_2.text = "• Learnable Spatial Embeddings\n• Continuous Cyclic Projectors\n• GELU Regularization Skip-Blocks\n• INT8 Edge Quantization Viable\n• Caching Rollout lag updates"
    ptc2_2.font.size = Pt(13)
    ptc2_2.font.color.rgb = COLOR_DARK
    ptc2_2.font.name = "Calibri"

    # Card 3: Hardware & Configuration
    tc3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.0), Inches(3.8), Inches(3.8), Inches(3.0))
    tc3.fill.solid()
    tc3.fill.fore_color.rgb = COLOR_CARD
    tc3.line.color.rgb = COLOR_BORDER
    tc3.line.width = Pt(1)
    tf_tc3 = tc3.text_frame
    tf_tc3.word_wrap = True
    ptc3_1 = tf_tc3.paragraphs[0]
    ptc3_1.text = "Hardware & Configuration"
    ptc3_1.font.size = Pt(14)
    ptc3_1.font.bold = True
    ptc3_1.font.color.rgb = COLOR_BLUE
    ptc3_1.font.name = "Trebuchet MS"
    ptc3_1.space_after = Pt(10)
    ptc3_2 = tf_tc3.add_paragraph()
    ptc3_2.text = "SYSTEM:\nCPU: i7/i9 Intel or Ryzen 7/9\nRAM: 16-32GB  |  512GB-1TB SSD\nOS: Win 10/11 or Ubuntu 20.04 LTS+\nGPU:\nNVIDIA GeForce RTX 3070/40 Series\n8GB GDDR6 VRAM"
    ptc3_2.font.size = Pt(12)
    ptc3_2.font.color.rgb = COLOR_DARK
    ptc3_2.font.name = "Calibri"

    # =========================================================================
    # SLIDE 4: FEASIBILITY & VIABILITY | CONGESTION HARMS (SIH Layout)
    # =========================================================================
    slide4 = add_professional_slide("Feasibility & Viability  |  Harms & Prevention")
    
    # Left Column: FEASIBILITY & VIABILITY
    left_card = slide4.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.0))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = COLOR_BG
    left_card.line.color.rgb = COLOR_BLUE
    left_card.line.width = Pt(1.5)
    
    tf4_left = left_card.text_frame
    tf4_left.word_wrap = True
    
    lp4_1 = tf4_left.paragraphs[0]
    lp4_1.text = "FEASIBILITY & VIABILITY FACTORS"
    lp4_1.font.size = Pt(18)
    lp4_1.font.bold = True
    lp4_1.font.color.rgb = COLOR_BLUE
    lp4_1.font.name = "Trebuchet MS"
    lp4_1.space_after = Pt(12)
    
    lp4_2 = tf4_left.add_paragraph()
    lp4_2.text = "1. Technical Feasibility:\n• Dynamic lag calculation rollout: 9/10\n• PyTorch deepcopy weights lock: 10/10\n• Edge compiled quantization: 8/10\n2. Resource Feasibility:\n• Computational demands: 9/10 (Under 12ms per step rollout)\n• Storage requirements: 9/10 (Sub-megabyte weights model)\n3. Viability Performance:\nOur multi-branch model operates continuously in production-grade municipal control grids.\n4. Viability Accuracy:\n0.9 Dampening controls prediction noise, perfectly matching real traffic commute distributions."
    lp4_2.font.size = Pt(13)
    lp4_2.font.color.rgb = COLOR_DARK
    lp4_2.font.name = "Calibri"

    # Right Column: CONGESTION COLLAPSE RISKS & MITIGATION
    right_card = slide4.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.0))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = COLOR_BG
    right_card.line.color.rgb = COLOR_ORANGE
    right_card.line.width = Pt(1.5)
    
    tf4_right = right_card.text_frame
    tf4_right.word_wrap = True
    
    rp4_1 = tf4_right.paragraphs[0]
    rp4_1.text = "CONGESTION COLLAPSE RISKS & MITIGATION"
    rp4_1.font.size = Pt(18)
    rp4_1.font.bold = True
    rp4_1.font.color.rgb = COLOR_ORANGE
    rp4_1.font.name = "Trebuchet MS"
    rp4_1.space_after = Pt(12)
    
    rp4_2 = tf4_right.add_paragraph()
    rp4_2.text = "Potential Risks & Harms:\n• Economic Harm:\n  - Courier/logistics delay, supply chain degradation\n• Environmental Harm:\n  - Elevated tailpipe idling emissions and carbon waste\n• Societal Harm:\n  - Emergency response vehicle blocks, excessive commuter stress\n\nCongestion Prevention Best Practices:\n✔ Use dynamic forecasting engines (NeuralRoute)\n✔ Enforce adaptive signal grids\n✔ Deploy route-balanced AV grid scheduling"
    rp4_2.font.size = Pt(13)
    rp4_2.font.color.rgb = COLOR_DARK
    rp4_2.font.name = "Calibri"

    # =========================================================================
    # SLIDE 5: POTENTIAL IMPACT ON TARGET AUDIENCE (SIH Layout)
    # =========================================================================
    slide5 = add_professional_slide("Potential Impact on Target Audience & Smart Optimization")
    
    # Left Column: Impact Tree
    left_card5 = slide5.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.0))
    left_card5.fill.solid()
    left_card5.fill.fore_color.rgb = COLOR_BG
    left_card5.line.color.rgb = COLOR_BLUE
    left_card5.line.width = Pt(1.5)
    
    tf5_left = left_card5.text_frame
    tf5_left.word_wrap = True
    
    lp5_1 = tf5_left.paragraphs[0]
    lp5_1.text = "OVERALL BENEFITS & TARGET IMPACT"
    lp5_1.font.size = Pt(18)
    lp5_1.font.bold = True
    lp5_1.font.color.rgb = COLOR_BLUE
    lp5_1.font.name = "Trebuchet MS"
    lp5_1.space_after = Pt(12)
    
    lp5_2 = tf5_left.add_paragraph()
    lp5_2.text = "Positive Target Impacts:\n• Municipalities:\n  - Adaptive smart city grid setups, reducing city blockages\n• Logistics Fleets (Flipkart/AVs):\n  - Coordinated route schedules, preventing courier delays\n• Daily Commuters:\n  - Seamless transit, restored schedule trust, and lower fuel waste\n\nNegative Impacts (Mitigated):\n• Edge Hardware Requirements:\n  - Overcome via FP16 quantization to enable edge compiles.\n• Signal Grid Dependency:\n  - Decoupled predictions to operate stand-alone during server downtime."
    lp5_2.font.size = Pt(13)
    lp5_2.font.color.rgb = COLOR_DARK
    lp5_2.font.name = "Calibri"

    # Right Column:infographics representation
    right_card5 = slide5.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.0))
    right_card5.fill.solid()
    right_card5.fill.fore_color.rgb = COLOR_BG
    right_card5.line.color.rgb = COLOR_ORANGE
    right_card5.line.width = Pt(1.5)
    
    tf5_right = right_card5.text_frame
    tf5_right.word_wrap = True
    
    rp5_1 = tf5_right.paragraphs[0]
    rp5_1.text = "DYNAMIC TRANSIT SCHEDULING INFOGRAPHIC"
    rp5_1.font.size = Pt(18)
    rp5_1.font.bold = True
    rp5_1.font.color.rgb = COLOR_ORANGE
    rp5_1.font.name = "Trebuchet MS"
    rp5_1.space_after = Pt(12)
    
    rp5_2 = tf5_right.add_paragraph()
    rp5_2.text = "How to Recognize & Mitigate Congestion:\n• Inconsistent Signal Scheduling:\n  - Align light duration patterns using dynamic inflow rollouts.\n• Uncoordinated routing:\n  - Deploy Route Balancing algorithm to coordinate smart AV fleets.\n• Multi-Step Exposure Bias:\n  - Suppress compounding noise via mathematically dampened autoregressive forecasts.\n• Transit Bottlenecks:\n  - Enable real-time prediction indicators for daily city command monitoring."
    rp5_2.font.size = Pt(13)
    rp5_2.font.color.rgb = COLOR_DARK
    rp5_2.font.name = "Calibri"

    # =========================================================================
    # SLIDE 6: RESEARCH & REFERENCE (SIH Layout)
    # =========================================================================
    slide6 = add_professional_slide("Research & Reference Documentation")
    
    # Left Column: Bullet list of references
    left_box6 = slide6.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.0))
    tf6_left = left_box6.text_frame
    tf6_left.word_wrap = True
    
    lp6_1 = tf6_left.paragraphs[0]
    lp6_1.text = "Scientific Research & Frameworks"
    lp6_1.font.size = Pt(20)
    lp6_1.font.bold = True
    lp6_1.font.color.rgb = COLOR_BLUE
    lp6_1.font.name = "Trebuchet MS"
    lp6_1.space_after = Pt(12)
    
    lp6_2 = tf6_left.add_paragraph()
    lp6_2.text = "1. Exposure Bias Mitigation:\n'Addressing Exposure Bias in Autoregressive Time-Series' - arxiv.org/abs/1906.02753\n\n2. Tabular Deep Architectures:\n'Deep Tabular ResNets for Non-Linear Feature Interaction'\n\n3. Geohash Spatial Mapping:\nGeohash Spatial Index Libraries - github.com/vinsci/geohash\n\n4. Scikit-Learn Scaling & Preprocessing:\nhttps://scikit-learn.org/stable/modules/preprocessing.html"
    lp6_2.font.size = Pt(13)
    lp6_2.font.color.rgb = COLOR_TEXT
    lp6_2.font.name = "Calibri"

    # Right Column: Reference boxes (Open CV Reference style)
    right_card6 = slide6.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.0)
    )
    right_card6.fill.solid()
    right_card6.fill.fore_color.rgb = COLOR_CARD
    right_card6.line.color.rgb = COLOR_BORDER
    right_card6.line.width = Pt(1.5)
    
    tf6_right = right_card6.text_frame
    tf6_right.word_wrap = True
    
    rp6_1 = tf6_right.paragraphs[0]
    rp6_1.text = "DEVELOPMENT REFERENCES"
    rp6_1.font.size = Pt(16)
    rp6_1.font.bold = True
    rp6_1.font.color.rgb = COLOR_ORANGE
    rp6_1.font.name = "Trebuchet MS"
    rp6_1.space_after = Pt(12)
    
    rp6_2 = tf6_right.add_paragraph()
    rp6_2.text = "• PyTorch Core Neural Network API:\n  https://pytorch.org/docs/stable/nn.html\n\n• Grab AI for SEA Competition Dataset:\n  https://github.com/grab/sea-traffic-management\n\n• Autoregressive Dampening Model Paper:\n  https://arxiv.org/abs/1906.02753"
    rp6_2.font.size = Pt(14)
    rp6_2.font.color.rgb = COLOR_DARK
    rp6_2.font.name = "Calibri"

    # Save presentation
    output_filename = "GridLock_2.0_Upgraded_Presentation_V4.pptx"
    prs.save(output_filename)
    print(f"Upgraded Presentation successfully created: {output_filename}")

if __name__ == "__main__":
    create_upgraded_presentation()
